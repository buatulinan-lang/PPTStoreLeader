"""Pembungkus grafik native PowerPoint dengan gaya template M-Flash."""
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_MARK
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from .theme import FONT, INK, MUTED, NAVY, SERIES_COLORS, LINE

TYPES = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
    "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "line": XL_CHART_TYPE.LINE,
    "line_marker": XL_CHART_TYPE.LINE_MARKERS,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "pie": XL_CHART_TYPE.PIE,
}


def _kosong(slide, x, y, w, h, pesan="Tidak ada data pada filter ini"):
    """Placeholder rapi bila grafik tidak punya data."""
    from .theme import rect, text, CARD, LINE, MUTED
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    rect(slide, x, y, w, h, fill=CARD, line=LINE)
    text(slide, x, y + h / 2 - 0.15, w, 0.3, pesan, 10, color=MUTED,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return None


def _font(obj, size=9, bold=False, color=MUTED):
    obj.font.name = FONT
    obj.font.size = Pt(size)
    obj.font.bold = bold
    obj.font.color.rgb = color


def add_chart(slide, kind, cats, series, x, y, w, h, colors=None, legend=False,
              labels=False, label_fmt=None, label_size=8, gap=60, overlap=None,
              cat_size=8.5, val_axis=True, cat_axis=True, smooth=False,
              line_width=2.0, hole=62, label_pos=None, label_color=None):
    cats = [str(x) for x in cats]
    if not cats or all((v is None) for vals in series.values() for v in vals):
        return _kosong(slide, x, y, w, h)
    cd = CategoryChartData()
    cd.categories = list(cats)
    for name, vals in series.items():
        cd.add_series(name, [None if v is None else float(v) for v in vals])
    gf = slide.shapes.add_chart(TYPES[kind], Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = gf.chart
    colors = colors or SERIES_COLORS

    ch.has_title = False
    ch.font.name = FONT
    ch.font.size = Pt(cat_size)
    ch.font.color.rgb = MUTED

    ch.has_legend = legend
    if legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        _font(ch.legend, 9, False, INK)

    plot = ch.plots[0]
    if kind in ("doughnut", "pie"):
        if kind == "doughnut":
            try:
                plot._element.get_or_add_holeSize().set("val", str(hole))
            except Exception:
                pass
        pts = plot.series[0].points
        for i, pt in enumerate(pts):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = colors[i % len(colors)]
            pt.format.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            pt.format.line.width = Pt(1.25)
    else:
        try:
            plot.gap_width = gap
            if overlap is not None:
                plot.overlap = overlap
        except Exception:
            pass
        for i, se in enumerate(ch.series):
            col = colors[i % len(colors)]
            if kind.startswith("line"):
                se.format.line.color.rgb = col
                se.format.line.width = Pt(line_width)
                se.smooth = smooth
            else:
                se.format.fill.solid()
                se.format.fill.fore_color.rgb = col
                se.format.line.fill.background()

    if labels:
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.name = FONT
        dl.font.size = Pt(label_size)
        dl.font.bold = True
        dl.font.color.rgb = label_color or INK
        if label_fmt:
            dl.number_format = label_fmt
            dl.number_format_is_linked = False
        if label_pos:
            try:
                dl.position = label_pos
            except Exception:
                pass

    if kind not in ("doughnut", "pie"):
        try:
            va = ch.value_axis
            va.has_major_gridlines = val_axis
            if val_axis:
                gl = va.major_gridlines.format.line
                gl.color.rgb = LINE
                gl.width = Pt(0.5)
            va.visible = val_axis
            va.has_minor_gridlines = False
            va.major_tick_mark = XL_TICK_MARK.NONE
            va.format.line.fill.background()
            _font(va.tick_labels, 8, False, MUTED)
        except Exception:
            pass
        try:
            ca = ch.category_axis
            ca.has_major_gridlines = False
            ca.visible = cat_axis
            ca.major_tick_mark = XL_TICK_MARK.NONE
            ca.format.line.color.rgb = LINE
            _font(ca.tick_labels, cat_size, False, MUTED)
        except Exception:
            pass
    return ch


def gauge(slide, x, y, w, h, value, color, hole=68):
    """Donat 2 segmen untuk kartu goal."""
    v = max(0.0, min(100.0, float(value)))
    ch = add_chart(slide, "doughnut", ["Tercapai", "Sisa"], {"s": [v, 100 - v]},
                   x, y, w, h, colors=[color, RGBColor(0xE6, 0xEA, 0xF2)], hole=hole)
    return ch
