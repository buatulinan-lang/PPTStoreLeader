"""Pembuat file PPTX 19 slide sesuai template standar M-Flash."""
from __future__ import annotations
import io
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from .theme import (base_slide, header, footer, kpi, note_card, rect, dot, text,
                    NAVY, NAVY2, INK, MUTED, GREEN, GREEN_D, RED, BLUE, AMBER,
                    CARD, CARD_ALT, WHITE, LINE, SW, SH, LOGO_MFLASH, LOGO_MADINAH, BG)
from .charts import add_chart, gauge
from . import metrics as M
from .metrics import n, pct, rp, tgl, tgl_s, periode_label
from .loader import STATUS_ORDER

STATUS_COLOR = {"Done": GREEN, "Cancel": RED, "Pending": AMBER, "Lainnya": BLUE}


def _web(slide):
    text(slide, 5.2, 6.81, 2.94, 0.35, "www.mflash.id", 14, color=NAVY, align=PP_ALIGN.CENTER)


# ============================================================ slide 1
def s_cover(prs, c):
    s = base_slide(prs, logos=False)
    s.shapes.add_picture(LOGO_MADINAH, Inches(1.35), Inches(1.75), Inches(2.25), Inches(1.5))
    s.shapes.add_picture(LOGO_MFLASH, Inches(1.45), Inches(3.75), Inches(2.16), Inches(1.9))
    text(s, 4.7, 2.35, 8.2, 0.57, c["judul"], 32, bold=True, color=NAVY)
    text(s, 4.7, 3.05, 8.2, 0.37, c["penyaji"], 15, color=MUTED)
    rect(s, 4.72, 3.82, 7.4, 0.04, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    text(s, 4.7, 4.12, 8.2, 0.4, c["lingkup"], 16, bold=True, color=INK)
    text(s, 4.7, 4.62, 8.2, 0.35, f"{n(c['r']['total'])} TRANSAKSI UNIK DIANALISIS", 14, color=MUTED)
    text(s, 4.7, 5.62, 8.2, 0.3, f"Dibuat otomatis dari dashboard · {c['dibuat']}", 11, color=MUTED)
    return s


# ============================================================ slide 2
def s_goal(prs, c):
    s = base_slide(prs, logos=False)
    text(s, 0.51, 0.52, 12.3, 0.8, c["goal_judul"], 40, bold=True, color=NAVY)
    text(s, 0.51, 1.35, 12.3, 0.4, c["goal_sub"], 16, color=MUTED)
    goals = c["goals"][:4]
    ng = max(1, len(goals))
    gap = 0.25
    cw = (12.1 - gap * (ng - 1)) / ng
    for i, g in enumerate(goals):
        x = 0.62 + i * (cw + gap)
        rect(s, x, 2.0, cw, 4.55, fill=CARD, line=LINE)
        text(s, x + 0.12, 2.26, cw - 0.24, 0.7, g["nama"], 15 if ng > 3 else 17,
             bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        v = float(g["nilai"])
        col = RED if v < 85 else (GREEN_D if v > 100 else AMBER)
        gw = min(2.9, cw - 0.5)
        gauge(s, x + (cw - gw) / 2, 3.0, gw, 2.5, min(v, 100), col)
        text(s, x + 0.12, 4.06, cw - 0.24, 0.5, pct(v, 2), 22 if ng > 3 else 24,
             bold=True, color=col, align=PP_ALIGN.CENTER)
        ket = g.get("ket") or (f"Aktual {pct(v,2)} · ditampilkan maks 100%" if v > 100
                               else f"Realisasi {pct(v,2)} dari target")
        text(s, x, 5.72, cw, 0.6, ket, 11, color=MUTED, align=PP_ALIGN.CENTER)
    for i, (lab, col) in enumerate([("< 85%", RED), ("85% - 100%", AMBER), ("> 100%", GREEN_D)]):
        dot(s, 0.6 + i * 1.55, 6.78, 0.16, col)
        text(s, 0.84 + i * 1.55, 6.75, 1.5, 0.28, lab, 12, color=MUTED)
    _web(s)
    return s


# ============================================================ slide 3
def s_catatan(prs, c):
    s = base_slide(prs)
    header(s, c["catatan_judul"], c["periode_label"])
    poin = [b for b in c["catatan"] if str(b).strip()][:6]
    rect(s, 0.62, 1.5, 12.1, max(1.2, 0.7 + 0.78 * len(poin)), fill=CARD, line=LINE)
    y = 1.85
    for i, b in enumerate(poin):
        dot(s, 1.0, y + 0.14, 0.14, NAVY)
        text(s, 1.35, y, 11.0, 0.6, b, 14, color=INK, spacing=1.15)
        y += 0.78
    footer(s, c["sumber"])
    return s


# ============================================================ slide 4
def s_ringkasan(prs, c):
    r, s = c["r"], base_slide(prs)
    header(s, "RINGKASAN KINERJA", c["periode_label"])
    if not r["total"]:
        return s
    cards = [("TOTAL UNIT MASUK", n(r["total"]), "sesuai filter aktif", NAVY),
             ("SELESAI (DONE)", n(r["done"]), f"{pct(r['p_done'])} dari total", GREEN_D),
             ("BATAL (CANCEL)", n(r["cancel"]), f"{pct(r['p_cancel'])} dari total", RED),
             ("PENDING", n(r["pending"]), f"{pct(r['p_pending'])} dari total", AMBER),
             ("RATA-RATA / HARI", n(r["rata_hari"], 1), f"{n(r['hari_periode'])} hari periode", NAVY)]
    for i, (l, v, sub, col) in enumerate(cards):
        kpi(s, 0.62 + i * 2.45, 1.42, 2.29, 1.32, l, v, sub, col)
    pb = c["per_bulan_status"]
    cats = [periode_label(p).split(" ")[0][:3] + " " + periode_label(p).split(" ")[1][2:] for p in pb.index]
    add_chart(s, "column_stacked", cats,
              {k: pb[k].tolist() for k in STATUS_ORDER if pb[k].sum() > 0},
              0.62, 3.05, 7.5, 3.3, colors=[STATUS_COLOR[k] for k in STATUS_ORDER if pb[k].sum() > 0],
              legend=True, gap=55, overlap=100)
    text(s, 8.35, 3.12, 4.37, 0.32, "Catatan Utama", 11, bold=True, color=NAVY)
    y = 3.58
    for b in c["catatan_ringkas"][:4]:
        dot(s, 8.35, y + 0.07, 0.1, NAVY)
        text(s, 8.61, y, 4.11, 0.6, b, 10.5, color=INK, spacing=1.1)
        y += 0.78
    footer(s, c["sumber"])
    return s


# ============================================================ slide 5
def s_komposisi(prs, c):
    r, s = c["r"], base_slide(prs)
    header(s, "KOMPOSISI STATUS PENGERJAAN", c["periode_label"])
    vals = [("Done", r["done"], r["p_done"]), ("Cancel", r["cancel"], r["p_cancel"]),
            ("Pending", r["pending"], r["p_pending"]), ("Lainnya", r["lain"], r["p_lain"])]
    vals = [v for v in vals if v[1] > 0]
    add_chart(s, "doughnut", [f"{k}  {pct(p)}" for k, _, p in vals],
              {"Jumlah": [v for _, v, _ in vals]}, 0.62, 1.5, 6.1, 4.9,
              colors=[STATUS_COLOR[k] for k, _, _ in vals], legend=True, hole=58)
    text(s, 7.0, 1.52, 5.72, 0.34, "Rincian Angka", 11, bold=True, color=NAVY)
    y = 2.02
    for k, v, p in vals:
        rect(s, 7.0, y, 5.72, 0.62, fill=CARD, line=LINE)
        dot(s, 7.2, y + 0.25, 0.13, STATUS_COLOR[k])
        text(s, 7.46, y, 2.3, 0.62, k, 11, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        text(s, 9.7, y, 1.5, 0.62, n(v), 11, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        text(s, 11.3, y, 1.2, 0.62, pct(p), 11, color=MUTED, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        y += 0.72
    note_card(s, 7.0, max(y + 0.1, 5.02), 5.72, 1.05,
              "Pending adalah kondisi terkini, bukan akumulasi sepanjang periode. Unit lama yang sudah "
              "tuntas tidak lagi terhitung di sini.")
    footer(s, c["sumber"])
    return s


# ============================================================ slide 6
def s_mom(prs, c):
    m, s = c["mom"], base_slide(prs)
    if not m:
        header(s, "BULAN BERJALAN VS BULAN SEBELUMNYA", "Data tidak cukup untuk perbandingan")
        return s
    header(s, "BULAN BERJALAN VS BULAN SEBELUMNYA",
           f"Dibandingkan setara: tanggal 1–{m['cutoff']} pada kedua bulan "
           f"({m['prev_label']} vs {m['cur_label']})")

    def d(cur, prev):
        v = M.delta(cur, prev)
        return "" if v is None else f"  ({'+' if v >= 0 else ''}{pct(v)})"
    cards = [("UNIT MASUK 1–%d" % m["cutoff"], n(m["b"]["total"]),
              f"{m['prev_label']}: {n(m['a']['total'])}{d(m['b']['total'], m['a']['total'])}"),
             ("SELESAI (DONE)", n(m["b"]["done"]),
              f"{m['prev_label']}: {n(m['a']['done'])}{d(m['b']['done'], m['a']['done'])}"),
             ("BATAL (CANCEL)", n(m["b"]["cancel"]),
              f"{m['prev_label']}: {n(m['a']['cancel'])}{d(m['b']['cancel'], m['a']['cancel'])}"),
             ("RATA-RATA / HARI", n(m["b"]["rata"], 1), f"{m['prev_label']}: {n(m['a']['rata'],1)}"),
             ("OMZET", rp(m["omzet_b"]), f"{m['prev_label']}: {rp(m['omzet_a'])}{d(m['omzet_b'], m['omzet_a'])}"),
             ("LABA KOTOR", rp(m["laba_b"]), f"{m['prev_label']}: {rp(m['laba_a'])}{d(m['laba_b'], m['laba_a'])}")]
    for i, (l, v, sub) in enumerate(cards):
        kpi(s, 0.62 + i * 2.04, 1.5, 1.88, 1.35, l, v, sub, NAVY, value_size=20)
    add_chart(s, "column", ["Total", "Done", "Cancel"],
              {m["prev_label"]: [m["a"]["total"], m["a"]["done"], m["a"]["cancel"]],
               m["cur_label"]: [m["b"]["total"], m["b"]["done"], m["b"]["cancel"]]},
              0.62, 3.15, 6.6, 3.2, colors=[NAVY2, NAVY], legend=True, labels=True, gap=80)
    add_chart(s, "line", [str(i) for i in m["daily_cur"].index],
              {m["prev_label"]: m["daily_prev"].tolist(), m["cur_label"]: m["daily_cur"].tolist()},
              7.45, 3.15, 5.27, 3.2, colors=[MUTED, NAVY], legend=True)
    text(s, 0.62, 6.62, 12.1, 0.34,
         "Pending tidak ikut dibandingkan: angkanya kondisi terkini, bukan kejadian bulan tersebut.",
         10, color=MUTED)
    footer(s, f"{m['cur_label']} baru berjalan {m['cutoff']} hari, sehingga {m['prev_label']} ikut dipotong "
              f"sampai tanggal {m['cutoff']} agar setara. Bila laju bertahan, {m['cur_label']} diperkirakan "
              f"menutup di sekitar {n(m['proyeksi'])} unit.")
    return s


# ============================================================ slide 7
def s_harian(prs, c):
    r, s = c["r"], base_slide(prs)
    h = c["harian"]
    header(s, "REKAP UNIT MASUK HARIAN", c["periode_label"])
    if not len(h):
        return s
    hi, lo = h.idxmax(), h.idxmin()
    avg = h.mean()
    cards = [("HARI AKTIF", n(r["hari_aktif"]), "hari ada unit masuk"),
             ("RATA-RATA / HARI", n(r["rata_hari_aktif"], 1), "unit per hari aktif"),
             ("HARI TERTINGGI", n(h.max()), f"{tgl(hi)} · {M.HARI[hi.dayofweek]}"),
             ("HARI TERENDAH", n(h.min()), f"{tgl(lo)} · {M.HARI[lo.dayofweek]}"),
             ("SELISIH TERTINGGI", f"+{pct(h.max()/avg*100-100, 0)}", "di atas rata-rata harian")]
    for i, (l, v, sub) in enumerate(cards):
        kpi(s, 0.62 + i * 2.45, 1.42, 2.29, 1.3, l, v, sub, NAVY)
    lab = [tgl_s(d) if i == 0 or d.day == 1 else " " * (i + 1) for i, d in enumerate(h.index)]
    add_chart(s, "line", lab, {"Unit masuk": h.tolist()}, 0.62, 3.0, 8.4, 3.35,
              colors=[NAVY], cat_size=7)
    text(s, 9.25, 3.05, 3.47, 0.32, "Rata-rata per Hari", 11, bold=True, color=NAVY)
    wd = c["weekday"]
    mx = wd.max() or 1
    for i in range(7):
        y = 3.5 + i * 0.42
        text(s, 9.25, y, 0.95, 0.3, M.HARI[i], 10, color=INK)
        rect(s, 10.25, y + 0.08, 1.75, 0.16, fill=CARD_ALT)
        rect(s, 10.25, y + 0.08, max(0.03, 1.75 * wd[i] / mx), 0.16, fill=NAVY if wd[i] < mx else GREEN)
        text(s, 12.05, y, 0.67, 0.3, n(wd[i], 0), 10, bold=True, color=INK, align=PP_ALIGN.RIGHT)
    footer(s, c["sumber"])
    return s


# ============================================================ slide 8
def s_top_hari(prs, c):
    s = base_slide(prs)
    t, h = c["top_hari"], c["harian"]
    if not len(t) or not len(h):
        header(s, "HARI DENGAN UNIT MASUK TERTINGGI", "Tidak ada data pada filter ini")
        return s
    avg = h.mean()
    b = t.iloc[0]
    header(s, "HARI DENGAN UNIT MASUK TERTINGGI",
           f"Puncak tertinggi: {tgl(b['TANGGAL'])} ({b['HARI']}) dengan {n(b['JUMLAH'])} unit — "
           f"{pct(b['VS'], 0)} di atas rata-rata harian ({n(avg, 0)} unit).")
    text(s, 0.62, 1.5, 6.5, 0.32, "10 Tanggal Tertinggi", 11, bold=True, color=NAVY)
    cols = [(0.74, 0.4, "#"), (1.16, 1.85, "TANGGAL"), (3.04, 1.15, "HARI"),
            (4.24, 1.05, "JUMLAH"), (5.34, 1.55, "VS RATA-RATA")]
    for x, w, lab in cols:
        text(s, x, 1.9, w, 0.28, lab, 8.5, bold=True, color=MUTED,
             align=PP_ALIGN.RIGHT if lab in ("JUMLAH", "VS RATA-RATA") else PP_ALIGN.LEFT)
    for i, row in t.iterrows():
        y = 2.22 + i * 0.42
        if i % 2 == 0:
            rect(s, 0.62, y, 6.0, 0.4, fill=CARD, shape=MSO_SHAPE.RECTANGLE)
        vals = [(0.74, 0.4, str(i + 1), False, MUTED, PP_ALIGN.LEFT),
                (1.16, 1.85, tgl(row["TANGGAL"]), False, INK, PP_ALIGN.LEFT),
                (3.04, 1.15, row["HARI"], False, MUTED, PP_ALIGN.LEFT),
                (4.24, 1.05, n(row["JUMLAH"]), True, INK, PP_ALIGN.RIGHT),
                (5.34, 1.55, f"+{pct(row['VS'], 0)}", True, GREEN_D, PP_ALIGN.RIGHT)]
        for x, w, v, bo, col, al in vals:
            text(s, x, y + 0.02, w, 0.36, v, 10.5, bold=bo, color=col, align=al,
                 anchor=MSO_ANCHOR.MIDDLE)
    tt = t.sort_values("JUMLAH")
    add_chart(s, "bar", [f"{tgl_s(r['TANGGAL'])} ({r['HARI'][:3]})" for _, r in tt.iterrows()],
              {"Unit": tt["JUMLAH"].tolist()}, 7.0, 1.5, 5.72, 4.0, colors=[NAVY],
              labels=True, val_axis=False)
    wd = c["weekday"]
    note_card(s, 7.0, 5.66, 5.72, 0.98,
              f"{M.HARI[int(wd.idxmax())]} rata-rata tersibuk ({n(wd.max(),0)} unit/hari), "
              f"{pct(wd.max()/max(wd.min(),1)*100-100, 0)} di atas {M.HARI[int(wd.idxmin())]} yang paling sepi "
              f"({n(wd.min(),0)}).")
    footer(s, c["sumber"])
    return s


# ============================================================ slide 9
def s_per_dim(prs, c):
    s = base_slide(prs)
    g = c["per_dim"].head(11)
    header(s, f"KINERJA PER {c['dim_label']}", "Diurutkan dari unit masuk terbesar")
    cols = [(0.74, 2.6, "NAMA", PP_ALIGN.LEFT), (3.5, 1.2, "UNIT MASUK", PP_ALIGN.RIGHT),
            (4.85, 1.0, "% DONE", PP_ALIGN.RIGHT), (6.0, 1.0, "% CANCEL", PP_ALIGN.RIGHT)]
    for x, w, lab, al in cols:
        text(s, x, 1.44, w, 0.3, lab, 8.5, bold=True, color=MUTED, align=al)
    for i, (k, row) in enumerate(g.iterrows()):
        y = 1.77 + i * 0.45
        if i % 2 == 0:
            rect(s, 0.6, y, 6.6, 0.43, fill=CARD, shape=MSO_SHAPE.RECTANGLE)
        cc = RED if row["P_CANCEL"] > c["r"]["p_cancel"] else MUTED
        for x, w, v, bo, col, al in [
                (0.74, 2.6, str(k)[:26], False, INK, PP_ALIGN.LEFT),
                (3.5, 1.2, n(row["UNIT"]), True, INK, PP_ALIGN.RIGHT),
                (4.85, 1.0, pct(row["P_DONE"]), False, GREEN_D, PP_ALIGN.RIGHT),
                (6.0, 1.0, pct(row["P_CANCEL"]), True, cc, PP_ALIGN.RIGHT)]:
            text(s, x, y, w, 0.43, v, 10, bold=bo, color=col, align=al, anchor=MSO_ANCHOR.MIDDLE)
    gg = g.head(8).iloc[::-1]
    add_chart(s, "bar_stacked", [str(i)[:18] for i in gg.index],
              {k: gg[k].tolist() for k in ["Done", "Cancel", "Pending"]},
              7.4, 1.44, 5.32, 5.0, colors=[GREEN, RED, AMBER], legend=True, gap=45, overlap=100)
    footer(s, c["sumber"])
    return s


# ============================================================ slide 10-12
def s_status_detail(prs, c, status, judul, catatan, color):
    s = base_slide(prs)
    d = c["detail"][status]
    header(s, judul, f"{n(d['jumlah'])} unit — {pct(d['persen'])} dari total pada filter aktif")
    tt, kk = d["teknisi_top"], d["kerusakan_top"]
    if c["dim_label"].upper() in ("TEKNISI", "NAMA TEKNISI"):
        kt = d["kategori_top"]
        kartu3 = ("KATEGORI TERBANYAK", str(kt[0])[:20], f"{n(kt[1])} unit ({pct(kt[2])})", NAVY)
    else:
        dt_ = c["dim_top"][status]
        kartu3 = (f"{c['dim_label']} TERBANYAK", str(dt_[0])[:20],
                  f"{n(dt_[1])} unit ({pct(dt_[2])})", NAVY)
    cards = [("JUMLAH", n(d["jumlah"]), f"{pct(d['persen'])} dari total", color),
             ("TEKNISI TERBANYAK", str(tt[0])[:20], f"{n(tt[1])} unit ({pct(tt[2])})", NAVY),
             kartu3,
             ("KERUSAKAN TERBANYAK", str(kk[0])[:20], f"{n(kk[1])} unit ({pct(kk[2])})", NAVY),
             ("RATA-RATA / HARI", n(d["rata_hari"], 2), f"{n(d['hari_periode'])} hari periode", NAVY)]
    for i, (l, v, sub, col) in enumerate(cards):
        kpi(s, 0.62 + i * 2.45, 1.42, 2.29, 1.32, l, v, sub, col, value_size=15 if i else 25)
    t1 = d["teknisi"].head(8).iloc[::-1]
    add_chart(s, "bar", [str(i)[:20] for i in t1.index], {"Jumlah": t1.tolist()},
              0.62, 3.0, 5.9, 3.3, colors=[color], labels=True, val_axis=False)
    t2 = d["kerusakan"].head(8).iloc[::-1]
    add_chart(s, "bar", [str(i)[:20] for i in t2.index], {"Jumlah": t2.tolist()},
              6.8, 3.0, 5.9, 3.3, colors=[NAVY], labels=True, val_axis=False)
    text(s, 0.62, 6.5, 12.1, 0.4, catatan, 10.5, color=INK)
    footer(s, c["sumber"])
    return s


# ============================================================ slide 13
def s_penjualan(prs, c):
    s = base_slide(prs)
    j = c["jual"]
    if not j:
        header(s, "PENJUALAN — MODAL, OMZET & LABA", "File faktur penjualan belum diunggah")
        return s
    header(s, "PENJUALAN — MODAL, OMZET & LABA",
           f"{c['periode_label']} · {n(j['faktur'])} faktur · {n(j['unit'])} unit terjual")
    cards = [("OMZET (HARGA JUAL)", rp(j["omzet"]), f"{n(j['baris'])} baris", NAVY),
             ("MODAL (HARGA BELI)", rp(j["modal"]), f"{pct(j['p_modal'])} dari omzet", MUTED),
             ("LABA KOTOR", rp(j["laba"]), f"margin {pct(j['margin'])}", GREEN_D),
             ("RATA-RATA / FAKTUR", rp(j["rata_faktur"]), f"laba {rp(j['laba_faktur'])}/faktur", NAVY),
             ("LABA / UNIT", rp(j["laba_unit"]), f"dari {n(j['unit'])} unit", NAVY)]
    for i, (l, v, sub, col) in enumerate(cards):
        kpi(s, 0.62 + i * 2.45, 1.42, 2.29, 1.32, l, v, sub, col, value_size=19)
    g = j["per_kategori"].head(6).iloc[::-1]
    add_chart(s, "bar_stacked", [str(i)[:16] for i in g.index],
              {"Modal (jt)": (g["MODAL"] / 1e6).tolist(), "Laba (jt)": (g["LABA"] / 1e6).tolist()},
              0.62, 3.0, 6.4, 3.4, colors=[NAVY2, GREEN], legend=True, gap=50, overlap=100)
    text(s, 7.3, 3.0, 5.42, 0.3, "Margin per kategori", 11, bold=True, color=NAVY)
    for x, w, lab, al in [(7.4, 1.6, "KATEGORI", PP_ALIGN.LEFT), (9.0, 1.2, "OMZET", PP_ALIGN.RIGHT),
                          (10.3, 1.2, "LABA", PP_ALIGN.RIGHT), (11.55, 1.05, "MARGIN", PP_ALIGN.RIGHT)]:
        text(s, x, 3.36, w, 0.28, lab, 8.5, bold=True, color=MUTED, align=al)
    for i, (k, row) in enumerate(j["per_kategori"].head(6).iterrows()):
        y = 3.68 + i * 0.45
        if i % 2 == 0:
            rect(s, 7.3, y, 5.42, 0.43, fill=CARD, shape=MSO_SHAPE.RECTANGLE)
        for x, w, v, bo, col, al in [(7.4, 1.6, str(k)[:14], True, INK, PP_ALIGN.LEFT),
                                     (9.0, 1.2, rp(row["OMZET"]), False, INK, PP_ALIGN.RIGHT),
                                     (10.3, 1.2, rp(row["LABA"]), False, GREEN_D, PP_ALIGN.RIGHT),
                                     (11.55, 1.05, pct(row["MARGIN"]), True, INK, PP_ALIGN.RIGHT)]:
            text(s, x, y, w, 0.43, v, 9.5, bold=bo, color=col, align=al, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, "Modal diambil dari kolom HARGA BELI yang sudah berupa total per baris. Kategori JASA "
              "umumnya bermodal nol sehingga marginnya tampil mendekati 100% — biaya tenaga kerja belum dibebankan di sini.")
    return s


# ============================================================ slide penjualan harian & bulanan
def s_jual_tren(prs, c):
    s = base_slide(prs)
    j = c["jual"]
    if not j:
        header(s, "REKAP PENJUALAN HARIAN & BULANAN", "File faktur penjualan belum diunggah")
        return s
    header(s, "REKAP PENJUALAN HARIAN & BULANAN", c["periode_label"])
    h = j["harian"]
    pb = j["per_bulan"]
    hari = len(h)
    best = h["OMZET"].idxmax()
    cards = [("OMZET TOTAL", rp(j["omzet"]), f"{n(hari)} hari ada penjualan", NAVY),
             ("RATA-RATA / HARI", rp(j["omzet"] / hari if hari else 0), "omzet per hari aktif", NAVY),
             ("HARI TERTINGGI", rp(h["OMZET"].max()), tgl(best), GREEN_D),
             ("LABA / HARI", rp(j["laba"] / hari if hari else 0), f"margin {pct(j['margin'])}", GREEN_D),
             ("FAKTUR / HARI", n(j["faktur"] / hari if hari else 0, 1), f"total {n(j['faktur'])} faktur", NAVY)]
    for i, (l, v, sub, col) in enumerate(cards):
        kpi(s, 0.62 + i * 2.45, 1.42, 2.29, 1.32, l, v, sub, col, value_size=19)
    lab = [tgl_s(d) if i == 0 or d.day == 1 else " " * (i + 1) for i, d in enumerate(h.index)]
    add_chart(s, "line", lab, {"Omzet (jt)": (h["OMZET"] / 1e6).tolist()},
              0.62, 3.0, 7.6, 3.35, colors=[NAVY], cat_size=7)
    text(s, 8.5, 3.0, 4.22, 0.3, "Omzet & Laba per Bulan", 11, bold=True, color=NAVY)
    add_chart(s, "column", [periode_label(p)[:3] for p in pb.index],
              {"Omzet (jt)": (pb["OMZET"] / 1e6).tolist(), "Laba (jt)": (pb["LABA"] / 1e6).tolist()},
              8.5, 3.35, 4.22, 3.0, colors=[NAVY, GREEN], legend=True, gap=60)
    footer(s, f"Omzet harian dari kolom TOTAL HARGA pada file rincian faktur penjualan. "
              f"Hari tertinggi: {tgl(best)} dengan {rp(h['OMZET'].max())}.")
    return s


# ============================================================ slide breakdown kategori penjualan
def s_jual_kategori(prs, c):
    s = base_slide(prs)
    j = c["jual"]
    if not j or not len(j.get("per_kategori_jual", [])):
        header(s, "PENJUALAN PER KATEGORI", "Kolom KATEGORI PENJUALAN tidak ditemukan")
        return s
    g = j["per_kategori_jual"]
    header(s, "PENJUALAN PER KATEGORI",
           f"{c['periode_label']} · {n(len(g))} kategori · diurutkan dari omzet terbesar")
    top = g.index[0]
    cards = [("KATEGORI TERBESAR", str(top)[:20], f"{rp(g.loc[top,'OMZET'])} ({pct(g.loc[top,'OMZET']/j['omzet']*100)})", NAVY),
             ("OMZET TOTAL", rp(j["omzet"]), f"{n(j['faktur'])} faktur", NAVY),
             ("LABA KOTOR", rp(j["laba"]), f"margin {pct(j['margin'])}", GREEN_D),
             ("MARGIN TERTINGGI", str(g["MARGIN"].idxmax())[:18], pct(g["MARGIN"].max()), GREEN_D),
             ("MARGIN TERENDAH", str(g["MARGIN"].idxmin())[:18], pct(g["MARGIN"].min()), RED)]
    for i, (l, v, sub, col) in enumerate(cards):
        kpi(s, 0.62 + i * 2.45, 1.42, 2.29, 1.32, l, v, sub, col, value_size=15 if i in (0, 3, 4) else 19)
    gg = g.head(8).iloc[::-1]
    add_chart(s, "bar_stacked", [str(i)[:18] for i in gg.index],
              {"Modal": (gg["MODAL"] / 1e6).tolist(), "Laba": (gg["LABA"] / 1e6).tolist()},
              0.62, 3.0, 6.4, 3.4, colors=[NAVY2, GREEN], legend=True, gap=45, overlap=100)
    text(s, 7.3, 3.0, 5.42, 0.3, "Rincian per kategori penjualan", 11, bold=True, color=NAVY)
    for x, w, lab, al in [(7.4, 1.85, "KATEGORI", PP_ALIGN.LEFT), (9.25, 1.0, "FAKTUR", PP_ALIGN.RIGHT),
                          (10.25, 1.3, "OMZET", PP_ALIGN.RIGHT), (11.6, 1.0, "MARGIN", PP_ALIGN.RIGHT)]:
        text(s, x, 3.36, w, 0.28, lab, 8.5, bold=True, color=MUTED, align=al)
    for i, (k, row) in enumerate(g.head(7).iterrows()):
        y = 3.68 + i * 0.4
        if i % 2 == 0:
            rect(s, 7.3, y, 5.42, 0.38, fill=CARD, shape=MSO_SHAPE.RECTANGLE)
        for x, w, v, bo, col, al in [(7.4, 1.85, str(k)[:18], True, INK, PP_ALIGN.LEFT),
                                     (9.25, 1.0, n(row["FAKTUR"]), False, MUTED, PP_ALIGN.RIGHT),
                                     (10.25, 1.3, rp(row["OMZET"]), False, INK, PP_ALIGN.RIGHT),
                                     (11.6, 1.0, pct(row["MARGIN"]), True, GREEN_D, PP_ALIGN.RIGHT)]:
            text(s, x, y, w, 0.38, v, 9, bold=bo, color=col, align=al, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, c["sumber"])
    return s


# ============================================================ slide 14
def s_voucher(prs, c):
    s = base_slide(prs)
    v = c["voucher"]
    if not v:
        header(s, c["voucher_judul"], "Tidak ada baris voucher pada filter aktif")
        note_card(s, 0.62, 2.0, 12.1, 1.2,
                  "Ubah kata kunci voucher di panel kiri aplikasi bila nama barangnya berbeda.")
        return s
    header(s, c["voucher_judul"],
           f"{c['periode_label']} · {n(v['qty'])} voucher terjual")
    cards = [("VOUCHER TERJUAL", n(v["qty"]), f"{n(v['trx'])} transaksi", NAVY),
             ("OMZET", rp(v["omzet"]), f"rata-rata {rp(v['rata_harga'])}/voucher", NAVY),
             ("MODAL", rp(v["modal"]), f"{pct(v['p_modal'])} dari omzet", MUTED),
             ("LABA KOTOR", rp(v["laba"]), f"margin {pct(v['margin'])}", GREEN_D),
             ("RATA-RATA / HARI", n(v["rata_hari"], 1), f"{n(v['hari'])} hari ada penjualan", NAVY)]
    for i, (l, sv, sub, col) in enumerate(cards):
        kpi(s, 0.62 + i * 2.45, 1.42, 2.29, 1.32, l, sv, sub, col, value_size=19)
    g = v["per_dim"].head(10).iloc[::-1]
    add_chart(s, "bar", [str(i)[:20] for i in g.index], {"Voucher": g.tolist()},
              0.62, 3.0, 6.4, 3.4, colors=[NAVY], labels=True, val_axis=False)
    pb = v["per_bulan"]
    add_chart(s, "column", [periode_label(p)[:3] for p in pb.index], {"Voucher": pb.tolist()},
              7.3, 3.0, 5.42, 3.4, colors=[BLUE], labels=True)
    footer(s, f"Terbanyak: {v['per_dim'].index[0]} ({n(v['per_dim'].iloc[0])} voucher). "
              f"Periode data: {tgl(v['mulai'])} – {tgl(v['akhir'])}.")
    return s


# ============================================================ slide 15
def s_bagi_hasil(prs, c):
    s = base_slide(prs)
    b = c["bagi"]
    if not b:
        header(s, "BAGI HASIL TEKNISI", "Data jasa tidak tersedia pada filter aktif")
        return s
    header(s, "BAGI HASIL TEKNISI", f"{c['periode_label']} · {n(b['teknisi'])} teknisi")
    cards = [("OMZET JASA", rp(b["omzet"]), f"{n(b['baris'])} baris", NAVY),
             ("BAGI HASIL (ATURAN)", rp(b["total"]), f"{pct(b['persen'])} dari omzet jasa", GREEN_D),
             (f"PEMBANDING FLAT {n(b['flat_rate'],0)}%", rp(b["flat"]), f"omzet jasa × {n(b['flat_rate'],0)}%", MUTED),
             ("SELISIH", rp(b["selisih"]), "flat lebih besar" if b["selisih"] < 0 else "aturan lebih besar",
              RED if b["selisih"] < 0 else GREEN_D),
             ("RATA-RATA / TEKNISI", rp(b["rata"]), f"dari {n(b['teknisi'])} teknisi", NAVY)]
    for i, (l, v, sub, col) in enumerate(cards):
        kpi(s, 0.62 + i * 2.45, 1.42, 2.29, 1.32, l, v, sub, col, value_size=19)
    g = b["per_teknisi"].head(10).iloc[::-1]
    add_chart(s, "bar", [str(i)[:22] for i in g.index], {"Bagi Hasil": (g["BAGI"] / 1e6).tolist()},
              0.62, 3.0, 6.4, 3.4, colors=[NAVY], labels=True, label_fmt='#,##0.0"jt"', val_axis=False)
    text(s, 7.3, 3.0, 5.42, 0.3, "Komposisi menurut tarif", 11, bold=True, color=NAVY)
    for x, w, lab, al in [(7.4, 1.5, "KATEGORI", PP_ALIGN.LEFT), (8.9, 0.9, "TARIF", PP_ALIGN.RIGHT),
                          (9.9, 1.3, "OMZET", PP_ALIGN.RIGHT), (11.3, 1.3, "BAGI HASIL", PP_ALIGN.RIGHT)]:
        text(s, x, 3.36, w, 0.28, lab, 8.5, bold=True, color=MUTED, align=al)
    for i, (k, row) in enumerate(b["per_kat"].head(5).iterrows()):
        y = 3.68 + i * 0.5
        if i % 2 == 0:
            rect(s, 7.3, y, 5.42, 0.48, fill=CARD, shape=MSO_SHAPE.RECTANGLE)
        for x, w, v, bo, col, al in [(7.4, 1.5, str(k).title(), True, INK, PP_ALIGN.LEFT),
                                     (8.9, 0.9, pct(row["TARIF"], 0), False, MUTED, PP_ALIGN.RIGHT),
                                     (9.9, 1.3, rp(row["OMZET"]), False, INK, PP_ALIGN.RIGHT),
                                     (11.3, 1.3, rp(row["BAGI"]), True, GREEN_D, PP_ALIGN.RIGHT)]:
            text(s, x, y, w, 0.48, v, 9.5, bold=bo, color=col, align=al, anchor=MSO_ANCHOR.MIDDLE)
    tar = " · ".join(f"{k.title()} {n(v,0)}%" for k, v in b["tarif"].items())
    footer(s, f"Tarif: {tar}. Prioritas bila dua kata kunci: Normal. Angka berbasis omzet jasa, "
              f"belum dikurangi biaya lain.")
    return s


# ============================================================ slide 16
def _tabel(slide, x, y, w, headers, rows, col_w=None, row_h=0.42, size=10.5):
    ncol = len(headers)
    col_w = col_w or [w / ncol] * ncol
    cx = x
    for i, hd in enumerate(headers):
        text(slide, cx + 0.1, y, col_w[i] - 0.2, 0.28, hd, 8.5, bold=True, color=MUTED)
        cx += col_w[i]
    yy = y + 0.32
    for r, row in enumerate(rows):
        if r % 2 == 0:
            rect(slide, x, yy, w, row_h, fill=CARD, shape=MSO_SHAPE.RECTANGLE)
        cx = x
        for i, v in enumerate(row):
            text(slide, cx + 0.1, yy, col_w[i] - 0.2, row_h, str(v), size, color=INK,
                 anchor=MSO_ANCHOR.MIDDLE)
            cx += col_w[i]
        yy += row_h
    return yy


def _kotak_jabatan(slide, x, y, w, h, nama, jabatan, fill, warna_teks=WHITE, size=11):
    rect(slide, x, y, w, h, fill=fill)
    text(slide, x + 0.1, y + 0.12, w - 0.2, 0.3, str(nama)[:30], size, bold=True,
         color=warna_teks, align=PP_ALIGN.CENTER)
    text(slide, x + 0.1, y + h - 0.4, w - 0.2, 0.3, str(jabatan)[:34], size - 2,
         color=warna_teks, align=PP_ALIGN.CENTER)


def _garis(slide, x1, y1, x2, y2):
    from pptx.util import Inches as _I
    cn = slide.shapes.add_connector(1, _I(x1), _I(y1), _I(x2), _I(y2))
    cn.line.color.rgb = LINE
    cn.line.width = Pt(1.25)
    return cn


def s_struktur(prs, c):
    s = base_slide(prs)
    header(s, "STRUKTUR ORGANISASI", c["lingkup"].title())
    org = c["struktur"]
    pembina, top, mid, bawah = org.get("pembina", []), org["pimpinan"], org["level2"], org["level3"]

    if not (pembina or top or mid or bawah):
        note_card(s, 0.62, 2.0, 12.1, 1.2,
                  "Isi tabel struktur organisasi (nama lengkap & jabatan) di aplikasi.")
        _web(s)
        return s

    y = 1.35 if pembina else 1.5

    # level 0 — ustadz pembina cabang
    if pembina:
        h0 = 0.8
        _kotak_jabatan(s, 5.17, y, 3.0, h0, pembina[0]["nama"], pembina[0]["jabatan"],
                       GREEN_D, size=12)
        _garis(s, 6.67, y + h0, 6.67, y + h0 + 0.3)
        y += h0 + 0.3

    # level 1 — pimpinan tertinggi cabang
    h1 = 0.85
    if top:
        _kotak_jabatan(s, 5.42, y, 2.5, h1, top[0]["nama"], top[0]["jabatan"], NAVY, size=12)
        _garis(s, 6.67, y + h1, 6.67, y + h1 + 0.32)
        y += h1 + 0.32

    # level 2 — supervisor & sejajar
    y2, h2 = y, 0.82
    n2 = max(1, len(mid))
    w2 = min(2.6, (12.1 - 0.3 * (n2 - 1)) / n2)
    total2 = n2 * w2 + 0.3 * (n2 - 1)
    x2 = (SW - total2) / 2
    pusat_spv = None
    if mid:
        _garis(s, x2 + w2 / 2, y2 - 0.16, x2 + total2 - w2 / 2, y2 - 0.16)
    for i, m in enumerate(mid):
        x = x2 + i * (w2 + 0.3)
        _kotak_jabatan(s, x, y2, w2, h2, m["nama"], m["jabatan"], NAVY2)
        _garis(s, x + w2 / 2, y2 - 0.16, x + w2 / 2, y2)
        if m.get("induk") and pusat_spv is None:
            pusat_spv = x + w2 / 2
    if mid:
        y = y2 + h2

    # level 3 — di bawah supervisor
    if bawah:
        y3 = y + 0.45
        per_baris = 5 if len(bawah) > 8 else 4
        w3, h3, gap = 2.3, 0.74, 0.22
        n_atas = min(per_baris, len(bawah))
        total_atas = n_atas * w3 + gap * (n_atas - 1)
        x_awal = (SW - total_atas) / 2
        bus = y3 - 0.22
        if pusat_spv:
            _garis(s, pusat_spv, y, pusat_spv, bus)
        if n_atas > 1:
            _garis(s, x_awal + w3 / 2, bus, x_awal + total_atas - w3 / 2, bus)
        for i, b in enumerate(bawah):
            baris, kol = divmod(i, per_baris)
            n_baris = min(per_baris, len(bawah) - baris * per_baris)
            total3 = n_baris * w3 + gap * (n_baris - 1)
            x = (SW - total3) / 2 + kol * (w3 + gap)
            yy = y3 + baris * (h3 + 0.22)
            if yy + h3 > 6.85:
                break
            if baris == 0:
                _garis(s, x + w3 / 2, bus, x + w3 / 2, yy)
            _kotak_jabatan(s, x, yy, w3, h3, b["nama"], b["jabatan"], CARD, INK, size=10)
            rect(s, x, yy, 0.06, h3, fill=BLUE, shape=MSO_SHAPE.RECTANGLE)
    _web(s)
    return s


def s_foto(prs, c, judul, kunci, subjudul=""):
    s = base_slide(prs)
    header(s, judul, subjudul or c["periode_label"])
    fotos = c.get(kunci) or []
    area = (0.62, 1.5, 12.1, 5.3)
    if not fotos:
        rect(s, *area, fill=CARD, line=LINE)
        text(s, area[0], area[1] + area[3] / 2 - 0.2, area[2], 0.4,
             "Unggah foto di aplikasi untuk mengisi slide ini", 13, color=MUTED,
             align=PP_ALIGN.CENTER)
        footer(s, c["sumber"])
        return s
    n_f = min(len(fotos), 4)
    kolom = 1 if n_f == 1 else 2
    baris = 1 if n_f <= 2 else 2
    gw = (area[2] - 0.3 * (kolom - 1)) / kolom
    gh = (area[3] - 0.3 * (baris - 1)) / baris
    for i, f in enumerate(fotos[:4]):
        r_, k_ = divmod(i, kolom)
        _pasang_foto(s, f, area[0] + k_ * (gw + 0.3), area[1] + r_ * (gh + 0.3), gw, gh)
    footer(s, c["sumber"])
    return s


def _pasang_foto(slide, foto, x, y, w, h):
    """Tempel foto proporsional di tengah kotak, dengan bingkai."""
    import io as _io
    from PIL import Image as _Image
    rect(slide, x, y, w, h, fill=CARD, line=LINE)
    src = _io.BytesIO(foto) if isinstance(foto, (bytes, bytearray)) else foto
    try:
        im = _Image.open(src)
        rasio = im.width / im.height
    except Exception:
        return
    if isinstance(src, _io.BytesIO):
        src.seek(0)
    pw, ph = w - 0.16, h - 0.16
    if rasio > pw / ph:
        fw, fh = pw, pw / rasio
    else:
        fh, fw = ph, ph * rasio
    slide.shapes.add_picture(src, Inches(x + (w - fw) / 2), Inches(y + (h - fh) / 2),
                             Inches(fw), Inches(fh))


# ============================================================ slide 17
def s_komitmen(prs, c):
    s = base_slide(prs)
    text(s, 0.51, 0.52, 11.18, 0.77, "2. KOMITMEN", 40, bold=True, color=NAVY)
    text(s, 0.92, 1.37, 8.21, 0.4, c["komitmen_intro"], 16, color=INK)
    rows = [[r.get("pencapaian", ""), r.get("komitmen", ""), r.get("target", "")]
            for r in c["komitmen"] if any(str(v).strip() for v in r.values())]
    if rows:
        rect(s, 0.65, 1.92, 12.02, min(4.8, 0.55 + 0.52 * len(rows)), fill=CARD, line=LINE)
        _tabel(s, 0.75, 2.05, 11.8, ["PENCAPAIAN", "KOMITMEN", "TARGET"],
               rows, col_w=[4.0, 4.4, 3.4], row_h=0.52)
    _web(s)
    return s


# ============================================================ slide 18
def s_kesimpulan(prs, c):
    s = base_slide(prs)
    header(s, "KESIMPULAN & TINDAK LANJUT", c["periode_label"])
    for i, (judul, isi) in enumerate(c["kesimpulan"][:4]):
        x = 0.62 + (i % 2) * 6.22
        y = 1.62 + (i // 2) * 2.42
        rect(s, x, y, 5.86, 2.06, fill=CARD, line=LINE)
        dot(s, x + 0.28, y + 0.3, 0.44, NAVY)
        text(s, x + 0.28, y + 0.3, 0.44, 0.44, str(i + 1), 15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + 0.9, y + 0.3, 4.7, 0.46, judul, 13, bold=True, color=NAVY)
        text(s, x + 0.3, y + 0.92, 5.28, 1.1, isi, 10.5, color=INK, spacing=1.15)
    footer(s, c["sumber"])
    return s


# ============================================================ slide 19
def s_penutup(prs, c):
    s = base_slide(prs, logos=False)
    rect(s, 0, 0, 5.39, 7.5, fill=RGBColor(0xF1, 0xF4, 0xFA), shape=MSO_SHAPE.RECTANGLE)
    s.shapes.add_picture(LOGO_MADINAH, Inches(1.05), Inches(1.59), Inches(2.53), Inches(1.73))
    s.shapes.add_picture(LOGO_MFLASH, Inches(1.26), Inches(4.18), Inches(2.1), Inches(1.87))
    text(s, 5.6, 3.12, 7.4, 1.06, "JAZAKUMULLAHU KHAIRAN", 40, bold=True, color=NAVY)
    text(s, 5.6, 4.2, 7.4, 0.4, c["penyaji"], 14, color=MUTED)
    return s


# ============================================================ build
def build(c) -> bytes:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
    s_cover(prs, c)
    s_goal(prs, c)
    s_catatan(prs, c)
    s_ringkasan(prs, c)
    s_komposisi(prs, c)
    s_mom(prs, c)
    s_harian(prs, c)
    s_top_hari(prs, c)
    s_per_dim(prs, c)
    s_status_detail(prs, c, "Pending", "UNIT TERTAHAN (PENDING)",
                    "Pending adalah beban kerja yang harus segera diurai — makin lama tertahan, "
                    "makin tinggi risiko komplain customer.", AMBER)
    s_status_detail(prs, c, "Done", "PENYELESAIAN (DONE)",
                    "Penyelesaian adalah indikator utama kapasitas cabang dan kecepatan teknisi.", GREEN_D)
    s_status_detail(prs, c, "Cancel", "PEMBATALAN (CANCEL)",
                    "Setiap pembatalan berpotensi menandakan persoalan di harga, waktu tunggu, "
                    "ketersediaan sparepart, atau komunikasi.", RED)
    s_penjualan(prs, c)
    s_jual_tren(prs, c)
    s_jual_kategori(prs, c)
    s_voucher(prs, c)
    s_bagi_hasil(prs, c)
    s_struktur(prs, c)
    s_foto(prs, c, "MEASURE ACTIVITY", "foto_measure")
    s_foto(prs, c, "AR", "foto_ar")
    s_komitmen(prs, c)
    s_kesimpulan(prs, c)
    s_penutup(prs, c)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
