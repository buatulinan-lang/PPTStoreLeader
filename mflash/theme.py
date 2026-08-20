"""Token desain yang diambil langsung dari template PPT M-Flash."""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

ASSETS = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "assets")
BG = os.path.join(ASSETS, "bg.jpg")
LOGO_MFLASH = os.path.join(ASSETS, "logo_mflash.png")
LOGO_MADINAH = os.path.join(ASSETS, "logo_madinah.png")

FONT = "Calibri"

NAVY = RGBColor(0x1F, 0x38, 0x64)
NAVY2 = RGBColor(0x2E, 0x53, 0x94)
INK = RGBColor(0x20, 0x24, 0x2E)
MUTED = RGBColor(0x6B, 0x72, 0x80)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
GREEN_D = RGBColor(0x1E, 0x88, 0x5A)
RED = RGBColor(0xC0, 0x39, 0x2B)
BLUE = RGBColor(0x2E, 0x9B, 0xD6)
AMBER = RGBColor(0xE1, 0x8C, 0x1F)
CARD = RGBColor(0xF6, 0xF8, 0xFC)
CARD_ALT = RGBColor(0xE6, 0xEA, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD5, 0xDD, 0xEB)

SERIES_COLORS = [NAVY, GREEN, RED, BLUE, AMBER, NAVY2]

SW, SH = 13.333, 7.5


def _tf(box, text, size, bold=False, color=INK, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, italic=False, spacing=None, wrap=True):
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = str(text).split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.line_spacing = spacing
        r = p.add_run()
        r.text = ln
        f = r.font
        f.name, f.size, f.bold, f.italic = FONT, Pt(size), bold, italic
        f.color.rgb = color
    return box


def text(slide, x, y, w, h, s, size=11, **kw):
    return _tf(slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)), s, size, **kw)


def rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.12, line_w=0.75):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    sp.text_frame.word_wrap = True
    return sp


def dot(slide, x, y, d, color):
    return rect(slide, x, y, d, d, fill=color, shape=MSO_SHAPE.OVAL)


def base_slide(prs, logos=True):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if os.path.exists(BG):
        s.shapes.add_picture(BG, 0, 0, Inches(SW), Inches(SH))
    if logos:
        if os.path.exists(LOGO_MFLASH):
            s.shapes.add_picture(LOGO_MFLASH, Inches(12.42), Inches(0.26), Inches(0.57), Inches(0.5))
        if os.path.exists(LOGO_MADINAH):
            s.shapes.add_picture(LOGO_MADINAH, Inches(11.4), Inches(0.32), Inches(0.57), Inches(0.38))
    return s


def header(slide, title, subtitle=""):
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(0.62), Inches(0.42),
                                 Inches(0.3), Inches(0.36))
    tri.rotation = 90
    tri.fill.solid(); tri.fill.fore_color.rgb = NAVY
    tri.line.fill.background(); tri.shadow.inherit = False
    text(slide, 1.06, 0.36, 11.6, 0.52, title, 28, bold=True, color=NAVY)
    if subtitle:
        text(slide, 1.06, 0.94, 11.6, 0.32, subtitle, 12, color=MUTED)


def footer(slide, note):
    text(slide, 0.62, 7.02, 12.1, 0.26, note, 9, color=MUTED)


def kpi(slide, x, y, w, h, label, value, sub, value_color=NAVY, value_size=25):
    rect(slide, x, y, w, h, fill=CARD, line=LINE)
    pad = 0.18
    text(slide, x + pad, y + 0.13, w - 2 * pad, 0.26, label, 8.5, bold=True, color=MUTED)
    text(slide, x + pad, y + 0.40, w - 2 * pad, 0.52, value, value_size, bold=True, color=value_color)
    text(slide, x + pad, y + h - 0.36, w - 2 * pad, 0.28, sub, 9, color=MUTED)


def note_card(slide, x, y, w, h, body, title=None):
    rect(slide, x, y, w, h, fill=CARD, line=LINE)
    ty = y + 0.1
    if title:
        text(slide, x + 0.2, ty, w - 0.4, 0.28, title, 10.5, bold=True, color=NAVY)
        ty += 0.34
    text(slide, x + 0.2, ty, w - 0.4, h - (ty - y) - 0.1, body, 10, color=INK, spacing=1.15)
